import io

from pptx import Presentation
from pypdf import PdfReader


def extract_text_from_pdf(data: bytes) -> str:
    reader = PdfReader(io.BytesIO(data))
    return "\n".join(page.extract_text() or "" for page in reader.pages)


def extract_text_from_pptx(data: bytes) -> str:
    prs = Presentation(io.BytesIO(data))
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text)
    return "\n".join(texts)


def extract_text(filename: str, data: bytes) -> str:
    lower = filename.lower()
    if lower.endswith(".pdf"):
        return extract_text_from_pdf(data)
    if lower.endswith(".pptx"):
        return extract_text_from_pptx(data)
    raise ValueError(f"Unsupported file type: {filename} (only .pdf and .pptx are supported)")
